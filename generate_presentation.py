#!/usr/bin/env python3
"""
지역경제동향 보고서 자동 생성 시스템 - 발표 PPT 생성기
20분 발표용 PowerPoint 파일을 자동 생성합니다.

사용법:
    python generate_presentation.py

출력:
    캡스톤_발표자료.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 설정
# ============================================================

# 색상 테마 (국가데이터처 느낌 - 파란색 계열)
COLORS = {
    'primary': RGBColor(0, 51, 102),      # 진한 파란색
    'secondary': RGBColor(0, 102, 153),   # 중간 파란색
    'accent': RGBColor(255, 193, 7),      # 노란색 (강조)
    'text_dark': RGBColor(33, 33, 33),    # 어두운 텍스트
    'text_light': RGBColor(255, 255, 255), # 흰색 텍스트
    'bg_light': RGBColor(240, 248, 255),  # 연한 파란 배경
    'success': RGBColor(76, 175, 80),     # 초록색
    'warning': RGBColor(255, 152, 0),     # 주황색
    'danger': RGBColor(244, 67, 54),      # 빨간색
}

# 폰트 설정
FONTS = {
    'title': '맑은 고딕',
    'body': '맑은 고딕',
    'code': 'Consolas',
}


# ============================================================
# 유틸리티 함수
# ============================================================

def add_title_slide(prs, title, subtitle=""):
    """표지/섹션 타이틀 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
    
    # 배경 색상 (그라데이션 효과를 위한 도형)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['primary']
    bg_shape.line.fill.background()
    
    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    p.alignment = PP_ALIGN.CENTER
    
    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.33), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_light']
        p.font.name = FONTS['body']
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, bullet_points, notes=""):
    """일반 콘텐츠 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 상단 색상 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    
    # 본문 내용
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(12), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 레벨 처리 (들여쓰기)
        if point.startswith('  - '):
            p.text = "    • " + point[4:]
            p.font.size = Pt(18)
            p.level = 1
        elif point.startswith('- '):
            p.text = "• " + point[2:]
            p.font.size = Pt(20)
        elif point == "":
            p.text = ""
            p.font.size = Pt(10)
        else:
            p.text = point
            p.font.size = Pt(20)
        
        p.font.color.rgb = COLORS['text_dark']
        p.font.name = FONTS['body']
        p.space_after = Pt(8)
    
    # 발표자 노트
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_two_column_slide(prs, title, left_content, right_content, notes=""):
    """2열 레이아웃 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 상단 색상 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    
    # 왼쪽 열
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(6), Inches(5.5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        p.font.name = FONTS['body']
        p.space_after = Pt(6)
    
    # 오른쪽 열
    right_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.5), Inches(6), Inches(5.5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        p.font.name = FONTS['body']
        p.space_after = Pt(6)
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_table_slide(prs, title, headers, rows, notes=""):
    """표가 포함된 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 상단 색상 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    
    # 표 생성
    cols = len(headers)
    row_count = len(rows) + 1  # 헤더 포함
    
    table_width = Inches(12)
    table_height = Inches(0.5 * row_count)
    left = Inches(0.67)
    top = Inches(1.8)
    
    table = slide.shapes.add_table(row_count, cols, left, top, table_width, table_height).table
    
    # 헤더 스타일
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = COLORS['text_light']
        para.alignment = PP_ALIGN.CENTER
    
    # 데이터 행
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            
            # 줄무늬 배경
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['bg_light']
            
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = COLORS['text_dark']
            para.alignment = PP_ALIGN.CENTER
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_highlight_box_slide(prs, title, boxes, notes=""):
    """강조 박스가 있는 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 상단 색상 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    
    # 박스들 배치
    box_width = Inches(3.8)
    box_height = Inches(2.2)
    start_left = Inches(0.7)
    start_top = Inches(1.8)
    gap = Inches(0.3)
    
    colors = [COLORS['secondary'], COLORS['success'], COLORS['warning'], COLORS['danger']]
    
    for i, (box_title, box_content) in enumerate(boxes):
        col = i % 3
        row = i // 3
        
        left = start_left + col * (box_width + gap)
        top = start_top + row * (box_height + gap)
        
        # 박스 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()
        
        # 박스 제목
        title_shape = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.1), box_width - Inches(0.2), Inches(0.5)
        )
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = box_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER
        
        # 박스 내용
        content_shape = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.6), box_width - Inches(0.2), box_height - Inches(0.7)
        )
        tf = content_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = box_content
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


# ============================================================
# 메인 PPT 생성
# ============================================================

def create_presentation():
    """발표 PPT 생성"""
    
    # 프레젠테이션 생성 (16:9 비율)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========================================
    # 슬라이드 1: 표지
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['primary']
    bg_shape.line.fill.background()
    
    # 메인 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(12.33), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 지역경제동향 보고서"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    p.alignment = PP_ALIGN.CENTER
    
    # 부제목
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(12.33), Inches(1)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "자동 생성 시스템"
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    p.alignment = PP_ALIGN.CENTER
    
    # 기관명
    org_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(12.33), Inches(0.5)
    )
    tf = org_box.text_frame
    p = tf.paragraphs[0]
    p.text = "의뢰 기관: 국가데이터처"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['body']
    p.alignment = PP_ALIGN.CENTER
    
    # 날짜
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5)
    )
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2025년 12월"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['body']
    p.alignment = PP_ALIGN.CENTER
    
    # ========================================
    # 슬라이드 2: 목차
    # ========================================
    add_content_slide(prs, "📋 목차", [
        "1. 프로젝트 개요 및 배경",
        "",
        "2. 시스템 아키텍처",
        "",
        "3. 주요 기능 (50개 보고서)",
        "",
        "4. 기술적 성과",
        "",
        "5. 프로젝트 수행 중 어려운 점",
        "",
        "6. 데모 시연",
        "",
        "7. 향후 개선 과제 및 결론",
    ], notes="목차를 간단히 소개하고 넘어가세요. (30초)")
    
    # ========================================
    # 슬라이드 3: 섹션 1 - 프로젝트 개요
    # ========================================
    add_title_slide(prs, "1. 프로젝트 개요", "문제 정의 및 목표")
    
    # ========================================
    # 슬라이드 4: 문제 정의
    # ========================================
    add_content_slide(prs, "😰 현재 업무 현황 (AS-IS)", [
        "- 국가데이터처는 66종의 승인통계를 작성",
        "",
        "- 매분기 '지역경제동향' 보고서 발간",
        "  - 10개 경제지표 × 17개 시도 = 50+ 페이지",
        "",
        "- 현재 보고서 작성 프로세스:",
        "  - 수작업으로 표, 그래프, 설명 문구 작성",
        "  - 전체 인력 총동원하여 약 1주일 소요",
        "  - 수기 입력 시 오타 및 숫자 오류 발생",
        "",
        "- 동일한 패턴의 반복 작업 → 비효율 발생",
    ], notes="국가데이터처의 현황과 문제점을 설명합니다. 1주일이 걸린다는 점을 강조하세요.")
    
    # ========================================
    # 슬라이드 5: 프로젝트 목표
    # ========================================
    add_table_slide(prs, "🎯 프로젝트 목표 (TO-BE)", 
        ["구분", "AS-IS", "TO-BE"],
        [
            ["소요 시간", "1주일", "수 시간"],
            ["투입 인력", "전체 인력 총동원", "1~2명"],
            ["오류 가능성", "높음 (수작업)", "낮음 (자동화)"],
            ["일관성", "담당자별 차이", "템플릿 기반 통일"],
            ["확장성", "특정 보고서 한정", "타 통계에도 적용 가능"],
        ],
        notes="AS-IS와 TO-BE를 대비하여 프로젝트의 가치를 명확히 전달하세요."
    )
    
    # ========================================
    # 슬라이드 6: 프로젝트 정보
    # ========================================
    add_content_slide(prs, "📌 프로젝트 정보", [
        "- 프로젝트명: 지역경제동향 보고서 자동 생성 시스템",
        "",
        "- 의뢰 기관: 국가데이터처 (Ministry of Data and Statistics)",
        "",
        "- 개발 기간: 2025년 12월 11일 ~ 12월 26일 (약 2주)",
        "",
        "- 기술 스택:",
        "  - Backend: Python, Flask",
        "  - Template: Jinja2",
        "  - Data: Pandas, OpenPyXL",
        "",
        "- 총 커밋 수: 104개",
    ], notes="프로젝트 기본 정보를 소개합니다.")
    
    # ========================================
    # 슬라이드 7: 섹션 2 - 시스템 아키텍처
    # ========================================
    add_title_slide(prs, "2. 시스템 아키텍처", "전체 구조 및 데이터 흐름")
    
    # ========================================
    # 슬라이드 8: 시스템 구조
    # ========================================
    add_content_slide(prs, "🏗️ 시스템 구조", [
        "- 입력 데이터:",
        "  - 분석표 엑셀 파일 (.xlsx)",
        "  - 기초자료 수집표 (.xlsx)",
        "",
        "- 처리 엔진:",
        "  - Flask 웹 애플리케이션",
        "  - Generator 모듈 (데이터 추출)",
        "  - Jinja2 템플릿 (서식 적용)",
        "",
        "- 출력 결과:",
        "  - HTML (웹 미리보기)",
        "  - HWPX (한글 문서)",
        "  - Excel (분석표 자동생성)",
        "  - PDF (인쇄용)",
    ], notes="시스템의 입력-처리-출력 구조를 설명합니다.")
    
    # ========================================
    # 슬라이드 9: 데이터 흐름
    # ========================================
    add_content_slide(prs, "🔄 데이터 처리 파이프라인", [
        "분석표 엑셀 → pandas DataFrame → Generator (정제)",
        "→ Jinja2 Template (렌더링) → HTML/HWPX 출력",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "- 각 보고서마다:",
        "  - *_generator.py : 데이터 추출 로직",
        "  - *_template.html : 서식 정의",
        "  - *_schema.json : 데이터 구조 명세",
        "",
        "- 새로운 보고서 추가 시:",
        "  - Generator + Template + Schema 3개 파일만 추가하면 됨",
    ], notes="Generator-Template-Schema 3요소 아키텍처를 강조하세요.")
    
    # ========================================
    # 슬라이드 10: 섹션 3 - 주요 기능
    # ========================================
    add_title_slide(prs, "3. 주요 기능", "50개 보고서 자동 생성")
    
    # ========================================
    # 슬라이드 11: 생성 가능한 보고서
    # ========================================
    add_highlight_box_slide(prs, "📄 생성 가능한 보고서 (총 50개)", [
        ("📋 요약 보고서 (9개)", "표지, 일러두기, 목차\n인포그래픽\n지역경제동향 요약 5개"),
        ("📊 부문별 보고서 (10개)", "광공업생산, 서비스업생산\n소비동향, 건설동향\n수출, 수입, 물가\n고용률, 실업률, 인구이동"),
        ("🗺️ 시도별 보고서 (18개)", "17개 시도별 보고서\n+ 참고 GRDP"),
        ("📈 통계표 (13개)", "각 부문별 통계표\n+ GRDP + 부록"),
    ], notes="50개 보고서를 카테고리별로 설명합니다.")
    
    # ========================================
    # 슬라이드 12: 웹 대시보드
    # ========================================
    add_content_slide(prs, "🖥️ 웹 대시보드 기능", [
        "- 📁 드래그 앤 드롭 파일 업로드",
        "  - 분석표 엑셀 파일을 간편하게 업로드",
        "",
        "- 👁️ 실시간 미리보기",
        "  - 클릭 즉시 보고서 미리보기 표시",
        "",
        "- 📑 4개 탭 구조",
        "  - 요약 / 부문별 / 시도별 / 통계표",
        "",
        "- ⚠️ 결측치 시각화",
        "  - 노란색 형광펜으로 누락 데이터 표시",
        "",
        "- 📥 다중 포맷 내보내기",
        "  - HTML, HWPX, PDF 지원",
    ], notes="대시보드의 주요 기능을 설명합니다. 데모에서 실제로 보여줄 예정임을 언급하세요.")
    
    # ========================================
    # 슬라이드 13: 섹션 4 - 기술적 성과
    # ========================================
    add_title_slide(prs, "4. 기술적 성과", "주요 구현 내용")
    
    # ========================================
    # 슬라이드 14: 스키마 기반 데이터 추출
    # ========================================
    add_content_slide(prs, "📐 스키마 기반 데이터 추출", [
        "- JSON 스키마로 데이터 매핑 정의",
        "  - 엑셀 시트명, 행/열 위치, 데이터 타입 명세",
        "  - 총 39개 스키마 파일 작성",
        "",
        "- 업종명 매핑 자동화",
        "  - 긴 공식 명칭 → 보고서용 축약형 변환",
        "  - 예: '전자 부품, 컴퓨터...' → '반도체·전자부품'",
        "",
        "- 확장성 확보",
        "  - 새로운 보고서: 스키마만 작성하면 됨",
        "  - 분석표 구조 변경: 스키마 수정으로 대응",
    ], notes="스키마 기반 아키텍처의 장점을 강조하세요.")
    
    # ========================================
    # 슬라이드 15: 가중치 기반 순위
    # ========================================
    add_content_slide(prs, "🏆 가중치 기반 순위 시스템", [
        "- 17개 시도의 경제지표 실시간 분석",
        "",
        "- 증감률에 가중치를 적용하여 순위 자동 계산",
        "  - 증가 지역: 증감률 높은 순 정렬",
        "  - 감소 지역: 증감률 낮은 순 정렬",
        "",
        "- 상위/하위 지역 자동 하이라이트",
        "  - 보고서 설명 문구에 자동 반영",
        "",
        "- 기여도 기반 업종 정렬",
        "  - 지역별 주요 증가/감소 업종 자동 식별",
    ], notes="순위 시스템이 보고서 설명 문구 자동 생성에 어떻게 활용되는지 설명하세요.")
    
    # ========================================
    # 슬라이드 16: 개발 통계
    # ========================================
    add_table_slide(prs, "📊 개발 통계", 
        ["항목", "수량"],
        [
            ["총 Git 커밋", "104개"],
            ["HTML 템플릿", "71개"],
            ["Python 생성기", "15개"],
            ["JSON 스키마", "39개"],
            ["API 엔드포인트", "10+ 개"],
            ["지원 보고서", "50개"],
            ["개발 기간", "16일"],
        ],
        notes="숫자로 프로젝트 규모를 어필하세요."
    )
    
    # ========================================
    # 슬라이드 17: 섹션 5 - 어려운 점
    # ========================================
    add_title_slide(prs, "5. 어려운 점", "프로젝트 수행 중 직면한 도전")
    
    # ========================================
    # 슬라이드 18: 어려운 점 1 - 비공개 자료
    # ========================================
    add_content_slide(prs, "😰 어려움 1: 비공개 자료 접근 제한", [
        "📌 상황:",
        "- 국가데이터처의 분석표, 기초자료 등 핵심 데이터가 비공개",
        "- 외부에서 실제 데이터에 접근하기 어려운 환경",
        "",
        "💪 극복:",
        "- 담당자와 주기적인 미팅을 통해 자료 확보",
        "- 제공받은 샘플 데이터를 기반으로 개발 진행",
        "- 정답 이미지(correct_answer/)를 참고하여 출력 형식 파악",
        "",
        "✅ 결과:",
        "- 제한된 자료 환경에서도 시스템 개발 완료",
        "- 실무 환경의 보안/비공개 정책 대응 경험 확보",
    ], notes="비공개 자료 환경에서 담당자와의 협업으로 극복한 점을 강조하세요.")
    
    # ========================================
    # 슬라이드 19: 어려운 점 2 - 팀원 중도포기
    # ========================================
    add_content_slide(prs, "😰 어려움 2: 팀원 중도포기", [
        "📌 상황:",
        "- 프로젝트 초기 3인 팀으로 시작",
        "- 1차 중간발표 직전/직후 팀원 2명 중도포기",
        "- 이후 1인 프로젝트로 전환",
        "",
        "💪 극복:",
        "- 핵심 기능 우선순위 재정립 후 집중 개발",
        "- Flask + Jinja2로 풀스택 통합 개발",
        "- 프론트엔드, 백엔드, 데이터 처리 전 영역 단독 수행",
        "",
        "✅ 결과:",
        "- 1인 풀스택 개발로 프로젝트 완수",
        "- 위기 대응 능력 및 자기주도적 문제해결 역량 향상",
    ], notes="팀원 이탈에도 불구하고 프로젝트를 완수한 점을 강조하세요.")
    
    # ========================================
    # 슬라이드 20: 어려운 점 3 - 서식 변경
    # ========================================
    add_content_slide(prs, "😰 어려움 3: 보고서 서식 변경", [
        "📌 상황:",
        "- 프로젝트 진행 중 2024년 4분기부터 보고서 서식 변경",
        "- 기존 개발한 템플릿 구조 재설계 필요",
        "",
        "💪 극복:",
        "- 템플릿 기반 모듈화 설계로 수정 범위 최소화",
        "- 데이터와 표현 계층 분리로 유지보수성 확보",
        "- 변경된 실제 보고서를 분석하여 신규 서식 파악",
        "",
        "✅ 결과:",
        "- 서식 변경에 성공적으로 대응",
        "- 확장 가능한 템플릿 구조 구축",
        "- 실무 환경과 동일한 불확실성 대응 경험",
    ], notes="외부 요인(서식 변경)에 유연하게 대응한 점을 강조하세요.")
    
    # ========================================
    # 슬라이드 21: 어려움 요약
    # ========================================
    add_table_slide(prs, "📋 어려움 극복 요약", 
        ["구분", "어려움", "극복", "강조 역량"],
        [
            ["데이터 접근", "비공개 자료", "담당자 협업, 샘플 활용", "커뮤니케이션, 문제해결"],
            ["팀 구성", "3인→1인 전환", "풀스택 단독 개발", "위기 대응, 자기주도성"],
            ["외부 요인", "서식 변경", "모듈화 템플릿 설계", "유연한 설계, 적응력"],
        ],
        notes="표로 간결하게 정리하여 마무리합니다."
    )
    
    # ========================================
    # 슬라이드 22: 섹션 6 - 데모
    # ========================================
    add_title_slide(prs, "6. 데모 시연", "실제 동작 확인")
    
    # ========================================
    # 슬라이드 23: 데모 안내
    # ========================================
    add_content_slide(prs, "🖥️ 데모 시연 내용", [
        "1️⃣ 대시보드 접속",
        "   - http://localhost:5050",
        "",
        "2️⃣ 분석표 엑셀 파일 업로드",
        "   - 드래그 앤 드롭으로 간편 업로드",
        "   - 연도/분기 자동 감지",
        "",
        "3️⃣ 보고서 미리보기",
        "   - 요약 → 부문별 → 시도별 → 통계표 순회",
        "   - 결측치 시각화 확인",
        "",
        "4️⃣ 보고서 내보내기",
        "   - HTML / HWPX 다운로드",
        "",
        "⏱️ 데모 시간: 약 3분",
    ], notes="데모를 시작하기 전에 보여줄 내용을 미리 안내합니다.")
    
    # ========================================
    # 슬라이드 24: 섹션 7 - 향후 계획
    # ========================================
    add_title_slide(prs, "7. 향후 개선 과제", "고도화 방향")
    
    # ========================================
    # 슬라이드 25: 향후 과제
    # ========================================
    add_content_slide(prs, "🚀 향후 개선 과제", [
        "🔴 높음 (단기):",
        "  - KOSIS API 연동: 실시간 통계 데이터 자동 수집",
        "  - 데이터 검증 자동화: JSON Schema 기반 유효성 검사",
        "",
        "🟡 중간 (중기):",
        "  - 차트 자동 생성: matplotlib/Chart.js 연동",
        "  - 버전 관리: 보고서 히스토리 추적",
        "  - 워크플로우 자동화: 승인 프로세스 구현",
        "",
        "🟢 낮음 (장기):",
        "  - LLM 활용: AI 기반 분석 문구 자동 생성",
        "  - 다른 보도자료 확장: 인구동향, 물가동향 등",
        "  - 클라우드 배포: Docker + Kubernetes",
    ], notes="향후 개선 과제를 우선순위별로 설명합니다.")
    
    # ========================================
    # 슬라이드 26: 기대 효과
    # ========================================
    add_highlight_box_slide(prs, "✨ 기대 효과", [
        ("⏱️ 시간 절감", "1주일 → 수 시간\n약 90% 이상 단축"),
        ("👥 인력 효율화", "반복 작업에서 해방\n고부가가치 업무 집중"),
        ("✅ 품질 향상", "휴먼 에러 최소화\n일관된 서식 적용"),
        ("🔧 확장 가능", "타 보도자료에도\n적용 가능한 프레임워크"),
    ], notes="프로젝트의 기대 효과를 강조합니다.")
    
    # ========================================
    # 슬라이드 27: 결론
    # ========================================
    add_content_slide(prs, "📝 결론", [
        "- 지역경제동향 보고서 자동 생성 시스템 구축 완료",
        "",
        "- 50개 보고서 자동 생성 지원",
        "",
        "- Generator-Template-Schema 아키텍처로 확장성 확보",
        "",
        "- 비공개 자료, 팀원 중도포기, 서식 변경 등 어려움 극복",
        "",
        "- 국가데이터처 업무 효율화에 기여할 것으로 기대",
        "",
        "",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "",
        "감사합니다. 질문 있으시면 말씀해 주세요! 🙏",
    ], notes="감사 인사와 함께 Q&A로 넘어갑니다.")
    
    # ========================================
    # 슬라이드 28: Q&A
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['primary']
    bg_shape.line.fill.background()
    
    # Q&A 텍스트
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Q & A"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['title']
    p.alignment = PP_ALIGN.CENTER
    
    # 부제목
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(12.33), Inches(1)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "질문 및 답변"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['text_light']
    p.font.name = FONTS['body']
    p.alignment = PP_ALIGN.CENTER
    
    # ========================================
    # 저장
    # ========================================
    output_path = "캡스톤_발표자료.pptx"
    prs.save(output_path)
    print(f"✅ PPT 파일 생성 완료: {output_path}")
    print(f"📊 총 슬라이드 수: {len(prs.slides)}개")
    
    return output_path


if __name__ == "__main__":
    create_presentation()

